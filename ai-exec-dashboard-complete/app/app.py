import os, io, json
import pandas as pd
import numpy as np
import streamlit as st
from datetime import datetime
from jinja2 import Template
from uuid import uuid4

# Optional exports
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt

st.set_page_config(page_title="AI Executive Dashboard", layout="wide")

@st.cache_data
def load_tables(data_dir):
    def read_csv(path):
        # Parse "month" if present
        head = open(path).read(1000)
        parse = ["month"] if "month" in head else []
        return pd.read_csv(path, parse_dates=parse) if parse else pd.read_csv(path)
    tables = {}
    for fname in os.listdir(data_dir):
        if fname.endswith(".csv"):
            tables[fname.replace(".csv","")] = read_csv(os.path.join(data_dir, fname))
    return tables

def mom(series: pd.Series):
    return (series - series.shift(1)) / series.shift(1)

def rolling3(series: pd.Series):
    return series.rolling(3).mean()

def first_of_month(dt):
    if isinstance(dt, pd.Timestamp):
        return dt.day == 1
    try:
        d = pd.to_datetime(dt)
        return d.day == 1
    except Exception:
        return False

def freshness_ok(df_map):
    latests = []
    for name, df in df_map.items():
        if "month" in df.columns:
            latests.append(df["month"].max())
    if not latests:
        return False, "No month columns present"
    latest = min(latests)  # ensure all tables are fresh enough
    age_days = (pd.Timestamp.today() - latest).days
    return age_days <= 45, f"Latest shared month {latest.date()} is {age_days} days old"

def reconcile_funding(kpi_funding, deals_n):
    m = kpi_funding["month"].max()
    a = float(kpi_funding.loc[kpi_funding["month"]==m, "total_usd"].iloc[0])
    b = float(deals_n.loc[deals_n["month"]==m, "amount_usd"].sum())
    if a == 0:
        return False, "kpi_funding total is zero"
    diff = abs(a-b)/a
    return diff <= 0.03, f"Reconciliation for {m.date()}: KPI={a:,.0f} vs Deals sum={b:,.0f} (Δ={diff*100:.2f}%)"

def outlier_flags(kdf, ici_df, jobs_df):
    flags = []
    def check(series, name):
        if len(series) < 2: return
        delta = mom(series).iloc[-1]
        if pd.notna(delta) and abs(delta) > 0.30:
            flags.append(f"{name} MoM change {delta*100:.1f}% > 30%")
    check(kdf["total_usd"], "Funding")
    check(ici_df["ici"], "ICI")
    check(jobs_df["postings_total"], "Job postings")
    return flags

def evidence_links(silver_df, m_col="month", url_col="source_url", month_val=None):
    if silver_df is None or url_col not in silver_df.columns: return "N/A"
    if month_val is None and m_col in silver_df.columns:
        month_val = silver_df[m_col].max()
    urls = (silver_df.loc[silver_df[m_col]==month_val, url_col]
            .dropna().drop_duplicates().head(5).tolist())
    return " • ".join(urls) if urls else "N/A"

def export_pdf(buf, month_str, metrics):
    c = canvas.Canvas(buf, pagesize=LETTER)
    width, height = LETTER
    y = height - 72
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, y, f"AI Executive Dashboard — {month_str}")
    y -= 24
    c.setFont("Helvetica", 11)
    for k, v in metrics.items():
        c.drawString(72, y, f"{k}: {v}")
        y -= 16
    c.showPage()
    c.save()
    return buf

def export_ppt(buf, month_str, metrics):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = f"AI Executive Dashboard — {month_str}"
    # bullets
    left = Inches(1); top = Inches(1.7); width = Inches(8); height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for k, v in metrics.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.level = 0
        p.font.size = Pt(14)
    prs.save(buf)
    return buf

# ---- App start
st.title("AI Executive Dashboard")
st.caption("Prototype with validation gates, executive scorecards, trends, leaderboards, and exports.")

data_dir = os.getenv("DATA_DIR", "data")
tables = load_tables(data_dir)

required = ["kpi_funding","kpi_adoption","kpi_inference_cost","kpi_benchmarks","kpi_oss","kpi_jobs","kpi_risk","deals_n"]
missing = [t for t in required if t not in tables]
if missing:
    st.error(f"Missing required tables: {missing}")
    st.stop()

# Basic schema checks
errors = []
if tables["kpi_funding"]["total_usd"].isnull().any(): errors.append("kpi_funding.total_usd has nulls")
if (tables["kpi_inference_cost"]["ici"] < 0).any(): errors.append("kpi_inference_cost.ici has negatives")
for name in required:
    df = tables[name]
    if "month" not in df.columns: errors.append(f"{name}.month column missing")
    elif not df["month"].apply(first_of_month).all():
        errors.append(f"{name}.month must be first-of-month dates")
fresh_ok, fresh_msg = freshness_ok({k:tables[k] for k in required})
if not fresh_ok: errors.append("Freshness failed: " + fresh_msg)
rec_ok, rec_msg = reconcile_funding(tables["kpi_funding"], tables["deals_n"])
if not rec_ok: errors.append("Reconciliation failed: " + rec_msg)

if errors:
    st.error("Validation failed — fix before publishing: " + " | ".join(errors))
else:
    st.success("Validation passed • " + fresh_msg + " • " + rec_msg)

lm = tables["kpi_funding"]["month"].max()
st.caption(f"Latest month: {lm.date()} • Data dir: {data_dir}")

# Scorecards
cols = st.columns(6)
k = tables["kpi_funding"].sort_values("month")
ici = tables["kpi_inference_cost"].sort_values("month")
ad = tables["kpi_adoption"].sort_values("month")
risk = tables["kpi_risk"].sort_values("month")
deals = tables["deals_n"]

with cols[0]:
    st.metric("Funding $", f"${k['total_usd'].iloc[-1]:,.0f}", delta=f"{mom(k['total_usd']).iloc[-1]*100:.1f}%")
with cols[1]:
    mna = int(deals.loc[deals["month"]==lm, "is_mna"].sum()) if "is_mna" in deals.columns else 0
    st.metric("M&A Count", f"{mna}")
with cols[2]:
    st.metric("Production Announcements", f"{ad['production_announcements'].iloc[-1]}", delta=f"{mom(ad['production_announcements']).iloc[-1]*100:.1f}%")
with cols[3]:
    st.metric("Inference Cost Index", f"{ici['ici'].iloc[-1]:.2f}", delta=f"{mom(ici['ici']).iloc[-1]*100:.1f}%")
with cols[4]:
    st.metric("Open-weights Share", f"{ici['open_weight_share'].iloc[-1]*100:.1f}%")
with cols[5]:
    st.metric("High-Severity Incidents", f"{risk['high_sev_incidents'].iloc[-1]}", delta=f"{mom(risk['high_sev_incidents']).iloc[-1]*100:.1f}%")

# Outlier flags
flags = outlier_flags(k, ici, tables["kpi_jobs"])
if flags:
    st.warning("Outlier flags: " + " | ".join(flags))

st.divider()

# Trends
import plotly.express as px
t1, t2 = st.columns(2)
with t1:
    fig = px.line(k, x="month", y=["total_usd","deal_count"], title="Funding $ and Deals (12 mo)")
    st.plotly_chart(fig, use_container_width=True)
with t2:
    ici2 = ici.copy()
    ici2["ici_ra3"] = rolling3(ici2["ici"])
    fig = px.line(ici2, x="month", y=["ici","ici_ra3"], title="Inference Cost Index (RA3)")
    st.plotly_chart(fig, use_container_width=True)

st.subheader("Leaderboards")
l1, l2 = st.columns(2)
with l1:
    top = deals[deals["month"]==lm].sort_values("amount_usd", ascending=False).head(10)
    st.write("Top 10 Deals")
    st.dataframe(top[["company","category","stage","amount_usd","source_url"]])
with l2:
    models = tables.get("models_n")
    if models is not None:
        latest_models = models[models["month"]==lm]
        st.write("Latest Model Releases")
        st.dataframe(latest_models[["model_name","family","open_weight","context_max","price_in_per_million","price_out_per_million","release_url"]])

st.subheader("Risk & Policy")
inc = tables.get("incidents_n")
if inc is not None:
    st.dataframe(inc[inc["month"]==lm][["org","incident_type","severity","summary","link"]])

# Evidence links (traceability)
st.caption("Evidence links (latest month): " + evidence_links(deals, month_val=lm))

# Executive Brief
brief_template = """
Headline: In {{ month_str }}, AI funding totaled ${{ funding | int | comma }}, while inference costs {{ 'fell' if ici_mom < 0 else 'rose' }} {{ (ici_mom*100)|round(1) }}% MoM.

3 Signals:
- Capital: Deals this month reached {{ deals }}.
- Capability: OSS momentum score at {{ oss | round(3) }}.
- Commercialization: Production announcements at {{ prods }}.

2 Risks:
- Policy uncertainty in key jurisdictions.
- Incident volume ticked up; review governance controls.

1 Decision:
- Focus pilots where cost-to-serve improves even if funding cools.
"""
class CommaTemplate(Template):
    def __init__(self, src):
        super().__init__(src)
        self.environment.filters["comma"] = lambda v: f"{v:,}"

tmpl = CommaTemplate(brief_template)
ici_mom = ((ici["ici"].iloc[-1] - ici["ici"].iloc[-2]) / ici["ici"].iloc[-2]) if len(ici) > 1 else 0
brief_text = tmpl.render(
    month_str=str(lm.date()),
    funding=int(k["total_usd"].iloc[-1]),
    deals=int(k["deal_count"].iloc[-1]),
    prods=int(ad["production_announcements"].iloc[-1]),
    oss=float(tables["kpi_oss"].sort_values("month")["oss_momentum_score"].iloc[-1]),
    ici_mom=float(ici_mom),
)

st.subheader("Executive Brief")
edited = st.text_area("Edit brief before export:", value=brief_text, height=220)

# Exports
colA, colB = st.columns(2)
with colA:
    if st.button("Export Executive Brief (PDF)"):
        from io import BytesIO
        buf = BytesIO()
        run_id = str(uuid4())[:8]
        metrics = {
            "Funding $": f"${k['total_usd'].iloc[-1]:,.0f}",
            "Deals": f"{int(k['deal_count'].iloc[-1])}",
            "ICI": f"{ici['ici'].iloc[-1]:.2f}",
            "Open-weights": f"{ici['open_weight_share'].iloc[-1]*100:.1f}%",
            "Production Announcements": f"{int(ad['production_announcements'].iloc[-1])}",
            "Incidents (High)": f"{int(risk['high_sev_incidents'].iloc[-1])}",
        }
        # write edited brief on first page
        from reportlab.lib.pagesizes import LETTER
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(buf, pagesize=LETTER)
        width, height = LETTER
        y = height - 72
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, y, f"Executive Brief — {str(lm.date())}")
        y -= 24
        c.setFont("Helvetica", 11)
        for line in edited.splitlines():
            c.drawString(72, y, line[:110])
            y -= 14
            if y < 72:
                c.showPage()
                y = height - 72
                c.setFont("Helvetica", 11)
        c.showPage()
        # metrics page
        y = height - 72
        c.setFont("Helvetica-Bold", 16); c.drawString(72, y, "Key Metrics"); y -= 24; c.setFont("Helvetica", 11)
        for kx, vx in metrics.items():
            c.drawString(72, y, f"{kx}: {vx}"); y -= 16
        c.showPage(); c.save()
        st.download_button("Download PDF", data=buf.getvalue(), file_name=f"ai_dashboard_{lm.date()}_{run_id}.pdf", mime="application/pdf")
with colB:
    if st.button("Export Key Slide (PPTX)"):
        from io import BytesIO
        buf = BytesIO()
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"AI Executive Dashboard — {str(lm.date())}"
        left = Inches(1); top = Inches(1.7); width = Inches(8); height = Inches(4.5)
        tx = slide.shapes.add_textbox(left, top, width, height).text_frame
        tx.word_wrap = True
        for line in edited.splitlines():
            p = tx.add_paragraph(); p.text = line; p.level = 0; p.font.size = Pt(14)
        prs.save(buf)
        run_id = str(uuid4())[:8]
        st.download_button("Download PPTX", data=buf.getvalue(), file_name=f"ai_dashboard_{lm.date()}_{run_id}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
